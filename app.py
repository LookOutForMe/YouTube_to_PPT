from flask import Flask, render_template, request, send_from_directory
from pytube import YouTube
import cv2
import pytesseract
from PIL import Image
from pptx import Presentation
import os

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'static/videos'
app.config['FRAME_FOLDER'] = 'static/frames'
app.config['OUTPUT_FOLDER'] = 'static/output'

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/convert', methods=['POST'])
def convert():
    video_url = request.form['video_url']
    youtube = YouTube(video_url)
    video_stream = youtube.streams.get_highest_resolution()
    video_path = os.path.join(app.config['UPLOAD_FOLDER'], 'video.mp4')
    video_stream.download(output_path=app.config['UPLOAD_FOLDER'], filename='video.mp4')

    os.makedirs(app.config['FRAME_FOLDER'], exist_ok=True)

    cap = cv2.VideoCapture(video_path)
    frame_count = 0
    ppt = Presentation()

    while cap.isOpened():
        ret, frame = cap.read()
        if not ret:
            break
        frame_count += 1
        image_path = os.path.join(app.config['FRAME_FOLDER'], f'frame_{frame_count}.png')
        cv2.imwrite(image_path, frame)

        image = Image.open(image_path)
        text = pytesseract.image_to_string(image)

        slide = ppt.slides.add_slide(ppt.slide_layouts[5])
        left = top = 0
        pic = slide.shapes.add_picture(image_path, left, top, height=ppt.slide_height)

        left_text = 0
        top_text = pic.top + pic.height + 10
        text_box = slide.shapes.add_textbox(left_text, top_text, pic.width, 100)
        text_frame = text_box.text_frame
        text_frame.text = text

    cap.release()

    output_pptx_path = os.path.join(app.config['OUTPUT_FOLDER'], 'output.pptx')
    ppt.save(output_pptx_path)

    return render_template('result.html', pptx_path=output_pptx_path)

@app.route('/download/<filename>')
def download_file(filename):
    return send_from_directory(app.config['OUTPUT_FOLDER'], filename)

if __name__ == '__main__':
    app.run(debug=True)
